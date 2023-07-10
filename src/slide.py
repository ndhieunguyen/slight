from pptx.util import Inches
from pptx import Presentation
import os
from bing_image_downloader import downloader
from .content import process_header
from PIL import Image
import shutil


def prepare_template(template_path):
    try:
        default_16_9_slide_size = (Inches(5.625), Inches(10))
        template = Presentation(template_path)
        if not (
            template.slide_height == default_16_9_slide_size[0] * 914400
            and template.slide_width == default_16_9_slide_size[1] * 914400
        ):
            pass
        else:
            template = Presentation()
            template.slide_height, template.slide_width = default_16_9_slide_size

    except:
        template = Presentation()
        template.slide_height, template.slide_width = default_16_9_slide_size

    return template


def get_layout_id(presentation):
    layout_id = 1
    while True:
        try:
            slide_layout = presentation.slide_layouts[layout_id]
        except:
            return None
        slide = presentation.slides.add_slide(slide_layout)
        try:
            shape = slide.placeholders[0]
        except:
            return None
        if shape.top / presentation.slide_height < 0.2:
            break
        layout_id += 1
    return layout_id


def create_slide(content, template, file_name, tmp_folder, font_file):
    layout_id = get_layout_id(template)
    for i in range(len(template.slides) - 1, -1, -1):
        rId = template.slides._sldIdLst[i].rId
        template.part.drop_rel(rId)
        del template.slides._sldIdLst[i]

    title = file_name
    slide_title = template.slides.add_slide(template.slide_layouts[0])
    title_box = slide_title.shapes.title
    title_box.text = title.upper()
    for i, place_holder in enumerate(slide_title.placeholders):
        if i < 1:
            continue
        sp = place_holder.element
        sp.getparent().remove(sp)

    key = list(content.keys())[0]

    img_slot = {"left": Inches(6), "top": Inches(1.6), "width": Inches(3.6), "height": Inches(3.6)}
    img_slot_ratio = img_slot["width"] / img_slot["height"]
    for item in content[key]:
        header, content = process_header(item["header"]), item["content"]
        image_query = f"{header}_{file_name}".replace(" ", "_")
        image = None
        if "Introduction" not in image_query:
            try:
                downloader.download(
                    image_query,
                    limit=1,
                    output_dir=tmp_folder,
                    force_replace=False,
                    adult_filter_off=False,
                    filter="photo",
                    timeout=10,
                    verbose=False,
                )

                image_folder_path = os.path.join(tmp_folder, image_query)
                image_path = os.path.join(image_folder_path, os.listdir(image_folder_path)[0])
                image = Image.open(image_path)

            except Exception as e:
                pass

        slide_layout = template.slide_layouts[layout_id]
        slide = template.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = header

        for i, place_holder in enumerate(slide.placeholders):
            if i < 1:
                continue
            sp = place_holder.element
            sp.getparent().remove(sp)

        if image:
            w, h = image.size
            try:
                if w / h > img_slot_ratio:  # image longer than slot -> resize image to fit slot_width
                    slide.shapes.add_picture(image_path, img_slot["left"], img_slot["top"], width=img_slot["width"])
                else:
                    slide.shapes.add_picture(image_path, img_slot["left"], img_slot["top"], height=img_slot["height"])
                shutil.rmtree(image_folder_path)
                del image

            except Exception as e:
                pass

        left, top, width, height = Inches(1), Inches(1.5), Inches(5), Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = content
        text_frame.fit_text(font_file=font_file, max_size=22)
        text_frame.word_wrap = True

    output_pptx_path = os.path.join(file_name.replace(" ", "_") + ".pptx")
    template.save(output_pptx_path)

    return output_pptx_path
